import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
import base64
from summa.summarizer import summarize as textrank_summarize
from pptx.dml.color import RGBColor
import os
# Function to upload and preview background image
def background_image_selector():
    st.sidebar.title("Background Image Selector")
    uploaded_bg_img = st.sidebar.file_uploader("Upload Background Image", type=["jpg", "jpeg", "png"])
    if uploaded_bg_img is not None:
        st.sidebar.image(uploaded_bg_img, caption="Uploaded Background Image", use_column_width=True)
    return uploaded_bg_img

# Function to extract text from PDF
def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# Function to extract text and images from PDF
def extract_text_and_images(pdf_file):
    text = ""
    images = []
    with fitz.open(pdf_file) as pdf:
        for page_num in range(len(pdf)):
            page = pdf.load_page(page_num)
            text += page.get_text()
            images += page.get_pixmap()

    return text, images

from pptx.util import Pt
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def generate_ppt(pdf_file, background_image, summary_length, right_m):
    prs = Presentation()
    reader = PdfReader(pdf_file)
    
    # Define maximum allowable width and height for the text box
    max_width = prs.slide_width - Inches(1)  # Subtracting a margin of 1 inch from each side
    max_height = prs.slide_height - Inches(1)  # Subtracting a margin of 1 inch from the top
    
    for page_num, page in enumerate(reader.pages):
        text = page.extract_text()
        #summarized_text = textrank_summarize(text)
        summarized_text = textrank_summarize(text, words=summary_length)

        # Add slide with blank layout
        slide_layout = prs.slide_layouts[1]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background color to white
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        
        # Add background image if available
        if background_image:
            pic = slide.shapes.add_picture(background_image, left=0, top=0, width=prs.slide_width, height=prs.slide_height)
            pic.lock_aspect_ratio = True  # Lock aspect ratio of the image
            pic.z_order = -1  # Set the z-order of the background image to the back

        # Add title
        title = slide.shapes.title
        title.text = f"Page {page_num + 1}"  # Set title to indicate page number
        
        # Calculate adjusted width and height for the text box
        width = Inches(3)  # Set the width to a maximum of 4 inches or the max_width, whichever is smaller
        height = min(Inches(5), max_height)  # Set the height to a maximum of 5 inches or the max_height, whichever is smaller
        
        # Add summary text as regular text on the slide
        left_margin = Inches(1)  # Set the left margin
        right_margin = Inches(2)  # Set the right margin
        available_width = prs.slide_width - left_margin - right_margin  # Calculate the available width for the textbox
        width =Inches(8)  # Set the width to a maximum of 4 inches or the available width, whichever is smaller

        # Calculate adjusted left margin based on available width and right margin
        left_margin = Inches(1)

        top_margin = Inches(right_m)   # Set the top margin
        txBox = slide.shapes.add_textbox(left_margin, top_margin, width, height)
        tf = txBox.text_frame
        tf.text = summarized_text
        p = tf.paragraphs[0]
        p.font.size = Pt(18)  # Set font size
        p.font.bold = True  # Optionally set font bold

    return prs






# Main Streamlit app
def main():
    st.title("PDF to Summarized PowerPoint Converter")
    
    # File upload
    pdf_file = st.file_uploader("Upload PDF file", type="pdf")
    # Background image selector
    uploaded_bg_img = background_image_selector()
    var=st.slider("top margin:", 1, 3, key="top_margin_slider")
    summary_length = st.slider("Number of Words for Summary per Slide", min_value=20, max_value=200, value=50, step=10, key="number_of_words")

    if pdf_file:
        st.subheader("Summarized Slides")
        ppt = generate_ppt(pdf_file,uploaded_bg_img, summary_length, var)
        
    # Load the PPT presentation
    prs = Presentation('output.pptx')
    
    # Create a Streamlit app
    st.title('PPT Slide Preview')
    
    # Iterate over the slides in the presentation
    for slide in prs.slides:
        # Initialize variables to store slide title and content
        title = ""
        content = ""
    
        # Iterate over shapes in the slide
        for shape in slide.shapes:
            # Check if the shape is a title
            if shape.shape_type == 1:  # Title placeholder
                title = shape.text
            # Check if the shape contains text
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    content += paragraph.text + "\n"
    
        # Display the slide title and content in Streamlit
        st.write(f'## {title}')
        st.write(content)
        st.markdown('<hr style="border-color:red; border-bottom: 5px solid rgba(21, 24, 57, 0.2);">', unsafe_allow_html=True)
        
        # Download button
    st.markdown(get_binary_file_downloader_html(ppt, 'Download PPT'), unsafe_allow_html=True)

# Function to create download link for PowerPoint
def get_binary_file_downloader_html(bin_file, file_label='File'):
    with open("output.pptx", 'wb') as f:
        bin_file.save(f)
    with open("output.pptx", 'rb') as f:
        data = f.read()
    bin_str = base64.b64encode(data).decode()
    href = f'<a href="data:application/octet-stream;base64,{bin_str}" download="{file_label}.pptx">Download {file_label}</a>'
    return href

if __name__ == "__main__":
    main()
