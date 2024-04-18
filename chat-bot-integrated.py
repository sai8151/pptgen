import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
import base64
from summa.summarizer import summarize as textrank_summarize
from pptx.dml.color import RGBColor
import os
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# Add imports for chatbot functionality
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
import spacy

# Function to upload and preview background image
def background_image_selector():
    st.sidebar.title("Background Image Selector")
    uploaded_bg_img = st.sidebar.file_uploader("Upload Background Image", type=["jpg", "jpeg", "png"])
    if uploaded_bg_img is not None:
        st.sidebar.image(uploaded_bg_img, caption="Uploaded Background Image", use_column_width=True)
    return uploaded_bg_img

# Function to extract text from PDF
def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
        
    return text

# Function to extract text and images from PDF
def extract_text_and_images(pdf_file):
    text = ""
    images = []
    with fitz.open(pdf_file) as pdf:
        for page_num in range(len(pdf)):
            page = pdf.load_page(page_num)
            text += page.get_text()
            images += page.get_pixmap()
    return text, images

# Function to preprocess text
def preprocess_text(text):
    tokens = word_tokenize(text.lower())
    stop_words = set(stopwords.words('english'))
    stemmer = PorterStemmer()
    preprocessed_text = " ".join([stemmer.stem(word) for word in tokens if word not in stop_words])
    return preprocessed_text

# Function to answer questions
def answer_question(text, question):
    # Preprocess question
    question_preprocessed = preprocess_text(question)

    # Tokenize text into sentences
    sentences = sent_tokenize(text)

    # Calculate relevance scores for each sentence
    relevance_scores = {}
    for i, sentence in enumerate(sentences):
        relevance_scores[i] = sum([1 for word in word_tokenize(sentence) if word in question_preprocessed])

    # Sort sentences based on relevance scores
    sorted_sentences = sorted(relevance_scores, key=relevance_scores.get, reverse=True)

    # Get top relevant sentences
    top_sentences_indices = sorted_sentences[:3]  # Adjust number of sentences as needed
    relevant_sentences = [sentences[i] for i in top_sentences_indices]

    # Concatenate relevant sentences to form the answer
    answer = ' '.join(relevant_sentences)
    
    return answer

# Function to generate PowerPoint slides
def generate_ppt(pdf_file, background_image, summary_length, top_m):
    prs = Presentation()
    reader = PdfReader(pdf_file)
    
    # Define maximum allowable width and height for the text box
    max_width = prs.slide_width - Inches(1)  # Subtracting a margin of 1 inch from each side
    max_height = prs.slide_height - Inches(1)  # Subtracting a margin of 1 inch from the top
    
    for page_num, page in enumerate(reader.pages):
        text = page.extract_text()
        summarized_text = textrank_summarize(text, words=summary_length)

        # Add slide with blank layout
        slide_layout = prs.slide_layouts[1]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background color to white
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        
        # Add background image if available
        if background_image:
            pic = slide.shapes.add_picture(background_image, left=0, top=0, width=prs.slide_width, height=prs.slide_height)
            pic.lock_aspect_ratio = True  # Lock aspect ratio of the image
            pic.z_order = -1  # Set the z-order of the background image to the back

        # Add title
        title = slide.shapes.title
        title.text = f"Page {page_num + 1}"  # Set title to indicate page number
        
        # Calculate adjusted width and height for the text box
        width = Inches(8.5)  # Set the width to a maximum of 4 inches or the max_width, whichever is smaller
        height = min(Inches(11), max_height)  # Set the height to a maximum of 5 inches or the max_height, whichever is smaller
        
        # Add summary text as regular text on the slide
        left_margin = Inches(0)  # Set the left margin
        right_margin = Inches(0)  # Set the right margin
        available_width = prs.slide_width - left_margin - right_margin  # Calculate the available width for the textbox
        width = Inches(8.5)  # Set the width to a maximum of 4 inches or the available width, whichever is smaller

        # Calculate adjusted left margin based on available width and right margin
        left_margin = Inches(0)

        top_margin = Inches(top_m)   # Set the top margin
        txBox = slide.shapes.add_textbox(left_margin , top_margin, width, height)
        tf = txBox.text_frame
        tf.text = summarized_text
        p = tf.paragraphs[0]
        p.font.size = Pt(18)  # Set font size
        p.font.bold = True  # Optionally set font bold

    return prs

# Main Streamlit app
def main():
    st.title("PDF to Summarized PowerPoint Converter")
    
    # File upload
    pdf_file = st.file_uploader("Upload PDF file", type="pdf")
    
    # Background image selector
    uploaded_bg_img = background_image_selector()
    
    # Slider for top margin
    var=st.slider("Top Margin:", 1, 3, key="top_margin_slider")
    
    # Slider for summary length
    summary_length = st.slider("Number of Words for Summary per Slide", min_value=20, max_value=200, value=50, step=10, key="number_of_words")

    # Chatbot section
    st.subheader("Chat with Bot")
    user_question = st.text_input("Ask me a question:", "")
    if st.button("Get Answer"):
        if pdf_file:
            text = extract_text_from_pdf(pdf_file)
            answer = answer_question(text, user_question)
            st.write(answer)

    # Generate PowerPoint slides
    if pdf_file:
        st.subheader("Summarized Slides")
        ppt = generate_ppt(pdf_file, uploaded_bg_img, summary_length, var)
        
    # Load and preview the generated PPT presentation
    prs = Presentation('output.pptx')
    
    # Display slide preview
    st.title('PPT Slide Preview')
    for slide in prs.slides:
        title = ""
        content = ""
        for shape in slide.shapes:
            if shape.shape_type == 1:  # Title placeholder
                title = shape.text
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    content += paragraph.text + "\n"
        st.write(f'## {title}')
        st.write(content)
        st.markdown('<hr style="border-color:red; border-bottom: 5px solid rgba(21, 24, 57, 0.2);">', unsafe_allow_html=True)
        
    # Download button for the PPT presentation
    st.markdown(get_binary_file_downloader_html(ppt, 'Download PPT'), unsafe_allow_html=True)

# Function to create download link for PowerPoint
def get_binary_file_downloader_html(bin_file, file_label='File'):
    with open("output.pptx", 'wb') as f:
        bin_file.save(f)
    with open("output.pptx", 'rb') as f:
        data = f.read()
    bin_str = base64.b64encode(data).decode()
    href = f'<a href="data:application/octet-stream;base64,{bin_str}" download="{file_label}.pptx">Download {file_label}</a>'
    return href

if __name__ == "__main__":
    main()
